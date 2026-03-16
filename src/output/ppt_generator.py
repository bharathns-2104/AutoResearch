"""
ppt_generator.py — Improved PPT Report Generator (Python wrapper)

Calls the Node.js PptxGenJS-based generator for high-quality output.
Falls back to the python-pptx approach if Node is unavailable.
"""

from __future__ import annotations

import json
import os
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict

OUTPUT_DIR = Path("reports")


class PPTGenerator:
    """
    Generates a professional PowerPoint report from mapped report data.
    Uses PptxGenJS (Node.js) for high-quality output.
    
    Usage:
        generator = PPTGenerator()
        generator.generate(mapped_data, output_path)
    """

    def __init__(self, output_dir: str = "reports"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        # Path to the JS generator script
        self._js_script = Path(__file__).parent / "ppt_generator.js"

    def generate(self, mapped_data: Dict[str, Any], file_path=None) -> str:
        """
        Generate the PPT report.
        
        Args:
            mapped_data: Output from ReportDataMapper.map()
            file_path:   Optional output path; auto-generated if None
        
        Returns:
            str: Path to the generated .pptx file
        """
        if not mapped_data:
            raise ValueError("mapped_data is empty — cannot generate PPT.")

        if file_path is None:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            file_path = self.output_dir / f"report_{ts}.pptx"

        file_path = str(file_path)

        # Try Node.js path first
        if self._node_available() and self._js_script.exists():
            return self._generate_via_node(mapped_data, file_path)

        # Fallback to python-pptx
        return self._generate_python_fallback(mapped_data, file_path)

    # ── Node.js path ──────────────────────────────────────────────────────────

    @staticmethod
    def _node_available() -> bool:
        try:
            result = subprocess.run(
                ["node", "--version"],
                capture_output=True, timeout=5
            )
            return result.returncode == 0
        except Exception:
            return False

    def _generate_via_node(self, mapped_data: Dict[str, Any], file_path: str) -> str:
        # Write mapped_data to a temp JSON file
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".json", delete=False, encoding="utf-8"
        ) as f:
            json.dump(mapped_data, f, default=str)
            tmp_json = f.name

        try:
            result = subprocess.run(
                ["node", str(self._js_script), tmp_json, file_path],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode != 0:
                raise RuntimeError(
                    f"Node PPT generator failed: {result.stderr[:500]}"
                )
            return file_path
        finally:
            os.unlink(tmp_json)

    # ── Python-pptx fallback ──────────────────────────────────────────────────

    def _generate_python_fallback(self, mapped: Dict[str, Any], file_path: str) -> str:
        """
        Clean python-pptx fallback — better than original but simpler than JS.
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        def rgb(hex_str: str):
            h = hex_str.lstrip("#")
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        NAVY  = "0D1B2A"
        TEAL  = "148F77"
        WHITE = "FFFFFF"

        prs = Presentation()
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(5.625)

        blank = prs.slide_layouts[6]

        def add_slide(title_text: str, dark: bool = False) -> any:
            slide = prs.slides.add_slide(blank)
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = rgb(NAVY if dark else "F4F6F7")
            # Title
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.15), Inches(9), Inches(0.7)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title_text
            run.font.size  = Pt(22)
            run.font.bold  = True
            run.font.color.rgb = rgb(WHITE if dark else NAVY)
            return slide

        def add_text_box(slide, text: str, x, y, w, h, size=11, color=None, bold=False):
            if color is None:
                color = WHITE
            tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, para_text in enumerate(text.split("\n")):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = para_text
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = rgb(color)

        # Slide 1: Cover
        s1 = add_slide(mapped.get("title_page", {}).get("project_title", "Report"), dark=True)
        so = mapped.get("score_overview", {})
        ds = mapped.get("domain_scores", {})
        dec = mapped.get("decision", {})
        score = so.get("overall_score", 0)
        add_text_box(s1,
            f"Overall Viability: {score:.0%}  ({so.get('rating','—')})\n"
            f"Financial: {ds.get('financial_score',0):.0%}  |  "
            f"Market: {ds.get('market_score',0):.0%}  |  "
            f"Competitive: {ds.get('competitive_score',0):.0%}\n"
            f"Decision: {dec.get('final_decision','—')}",
            0.5, 1.5, 9, 2.0, size=14, color=WHITE
        )

        # Slide 2: Executive Summary
        s2 = add_slide("Executive Summary")
        summary = (mapped.get("executive_summary") or {}).get("summary_text", "")
        add_text_box(s2, summary[:800], 0.5, 1.0, 9, 3.5, size=11, color="1C2833")

        # Slide 3: Scores
        s3 = add_slide("Viability Score Overview")
        lines = [
            f"Overall Viability Score: {so.get('overall_score',0):.0%}  ({so.get('rating','—')})",
            f"Financial Score:         {ds.get('financial_score',0):.0%}",
            f"Market Score:            {ds.get('market_score',0):.0%}",
            f"Competitive Score:       {ds.get('competitive_score',0):.0%}",
        ]
        add_text_box(s3, "\n".join(lines), 0.5, 1.1, 9, 3.0, size=13, color="1C2833")

        # Slide 4-6: Detail slides
        for section_key, section_title in [
            ("financial_details", "Financial Analysis"),
            ("market_details",    "Market Analysis"),
            ("competitive_details","Competitive Analysis"),
        ]:
            s = add_slide(section_title)
            data = mapped.get(section_key) or {}
            lines = [f"{k}: {v}" for k, v in list(data.items())[:8] if not isinstance(v, (dict, list))]
            add_text_box(s, "\n".join(lines), 0.5, 1.1, 9, 3.5, size=11, color="1C2833")

        # Slide 7: Risks
        s_risk = add_slide("Risk Assessment")
        risks = (mapped.get("risk_analysis") or {}).get("risks", [])
        risk_lines = [f"[{r.get('severity','?')}] {r.get('category','')}: {r.get('message','')}" for r in risks[:6]]
        add_text_box(s_risk, "\n".join(risk_lines) or "No risks identified.", 0.5, 1.1, 9, 3.5, size=10, color="1C2833")

        # Slide 8: Recommendations
        s_rec = add_slide("Strategic Recommendations")
        recs = (mapped.get("recommendations") or {}).get("recommendations", [])
        add_text_box(s_rec, "\n".join(f"{i+1}. {r}" for i, r in enumerate(recs[:6])), 0.5, 1.1, 9, 3.5, size=11, color="1C2833")

        # Slide 9: Decision (dark)
        s_dec = add_slide(f"Decision: {dec.get('final_decision','—')}", dark=True)
        add_text_box(s_dec,
            f"Viability Score: {so.get('overall_score',0):.0%}\n"
            f"Rating: {so.get('rating','—')}",
            0.5, 2.0, 9, 1.5, size=16, color=WHITE
        )

        prs.save(file_path)
        return file_path